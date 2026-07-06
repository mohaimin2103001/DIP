from __future__ import annotations

import argparse
from pathlib import Path

import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


PROJECT_ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = Path(__file__).resolve().parent


def render_pdf_pages(pdf_path: Path, output_dir: Path, zoom: float) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    for old_image in output_dir.glob("slide_*.png"):
        old_image.unlink()

    page_images: list[Path] = []
    document = fitz.open(pdf_path)
    matrix = fitz.Matrix(zoom, zoom)
    try:
        for page_number, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_path = output_dir / f"slide_{page_number:02d}.png"
            pixmap.save(image_path)
            page_images.append(image_path)
    finally:
        document.close()

    return page_images


def add_full_slide_image(prs: Presentation, image_path: Path) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    with Image.open(image_path) as image:
        image_width, image_height = image.size

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_ratio = slide_width / slide_height
    image_ratio = image_width / image_height

    if image_ratio > slide_ratio:
        height = slide_height
        width = int(height * image_ratio)
        left = int((slide_width - width) / 2)
        top = 0
    else:
        width = slide_width
        height = int(width / image_ratio)
        left = 0
        top = int((slide_height - height) / 2)

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def build_pptx(page_images: list[Path], pptx_path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)

    for image_path in page_images:
        add_full_slide_image(presentation, image_path)

    presentation.save(pptx_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert a PDF slide deck into an uneditable PPTX by placing each rendered page as a full-slide image."
    )
    parser.add_argument(
        "--pdf",
        type=Path,
        default=REPORT_DIR / "retinal_vessel_segmentation_slides.pdf",
        help="Input PDF slides file.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=REPORT_DIR / "retinal_vessel_segmentation_slides_uneditable_ULTRA.pptx",
        help="Output PPTX file.",
    )
    parser.add_argument(
        "--images",
        type=Path,
        default=REPORT_DIR / "slides_page_images_ultra",
        help="Folder where rendered slide PNG images are saved.",
    )
    parser.add_argument(
        "--zoom",
        type=float,
        default=10.0,
        help="PDF render scale. 10.0 gives about 4536 x 2552 px for this Beamer deck.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    pdf_path = args.pdf.resolve()
    output_path = args.output.resolve()
    image_dir = args.images.resolve()

    if not pdf_path.exists():
        raise FileNotFoundError(f"Missing input PDF: {pdf_path}")

    page_images = render_pdf_pages(pdf_path, image_dir, args.zoom)
    build_pptx(page_images, output_path)

    first_image_size = "unknown"
    if page_images:
        with Image.open(page_images[0]) as image:
            first_image_size = f"{image.size[0]} x {image.size[1]}"

    print(f"PDF pages converted: {len(page_images)}")
    print(f"Rendered image size: {first_image_size}")
    print(f"Rendered images: {image_dir}")
    print(f"PowerPoint created: {output_path}")
    print(f"PowerPoint size: {output_path.stat().st_size / (1024 * 1024):.2f} MB")


if __name__ == "__main__":
    main()
